"""Build a short Redwood-style PowerPoint for the OCI AI Agents demo."""

from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
TEMPLATE = Path(
    "/Users/liviarodrigues/Documents/Codex/2026-06-18/"
    "redwood-powerpoint-creator-plugin-redwood-powerpoint/work/presentations/"
    "data-ai-lad-team-demos/tmp/template-starter.pptx"
)
OUTPUT = DOCS / "tdc-oci-ai-agents-demo-oracle.pptx"
ARCH = DOCS / "arquitetura-oci-ai-agents-tdc.png"
QR = DOCS / "github-qr.png"
GITHUB_URL = "https://github.com/LiviaFernandes/tdc-oci-ai-agents-lab"


COLORS = {
    "bark": RGBColor(49, 45, 42),
    "muted": RGBColor(92, 86, 80),
    "sienna": RGBColor(174, 86, 44),
    "ocean": RGBColor(44, 89, 103),
    "ivy": RGBColor(117, 156, 108),
    "rose": RGBColor(163, 100, 114),
    "air": RGBColor(252, 251, 250),
    "neutral": RGBColor(245, 244, 242),
    "oracle": RGBColor(199, 70, 52),
    "white": RGBColor(255, 255, 255),
}


def remove_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[index]
    prs.part.drop_rel(slide_id.rId)
    del slide_id_list[index]


def delete_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def clear_prompt_text(slide):
    stale_terms = [
        "CX Title",
        "Subhead goes here",
        "Name",
        "Presenter",
        "Title on a light",
        "Subtitle",
        "Main paragraph",
        "First level bullet",
        "An end slide",
    ]
    for shape in list(slide.shapes):
        text = shape.text.strip() if hasattr(shape, "text") else ""
        if text and any(term in text for term in stale_terms):
            delete_shape(shape)


def set_shape_text(shape, text, size=6, color="muted"):
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Oracle Sans"
    run.font.size = Pt(size)
    run.font.color.rgb = COLORS[color]


def normalize_template_footer(slide, slide_number):
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = shape.text.strip()
        if not text:
            continue
        if "Copyright" in text:
            set_shape_text(
                shape,
                "Copyright © 2026, Oracle and/or its affiliates | Confidential: Internal",
                size=6,
                color="muted",
            )
        elif text.isdigit():
            set_shape_text(shape, str(slide_number), size=6, color="muted")


def textbox(slide, text, x, y, w, h, size=24, bold=False, color="bark", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Oracle Sans"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def bullets(slide, items, x, y, w, h, size=19, color="bark"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Oracle Sans"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(10)
    return box


def pill(slide, text, x, y, w, h, fill="neutral", line="sienna", size=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    shape.line.width = Pt(1.25)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Oracle Sans"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = COLORS["bark"]
    return shape


def add_link(slide, text, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.hyperlink.address = GITHUB_URL
    run.font.name = "Oracle Sans"
    run.font.size = Pt(16)
    run.font.color.rgb = COLORS["ocean"]
    return box


def build():
    DOCS.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(TEMPLATE, OUTPUT)
    prs = Presentation(str(OUTPUT))

    # Keep title, three content slides, and the closing slide.
    for index in [7, 6, 5, 4]:
        remove_slide(prs, index)

    slides = list(prs.slides)
    for slide in slides:
        clear_prompt_text(slide)
    for idx, slide in enumerate(slides, start=1):
        normalize_template_footer(slide, idx)

    # Slide 1
    s = slides[0]
    textbox(s, "Lab TDC: AI Agents na OCI", 0.75, 0.75, 7.4, 0.65, size=34, bold=True, color="white")
    textbox(s, "RAG, Custom Tools e endpoint em uma demo de 1 hora", 0.78, 1.42, 7.8, 0.35, size=18, color="white")
    textbox(s, "Material de apoio", 0.78, 4.55, 2.5, 0.3, size=14, color="white")
    add_link(s, GITHUB_URL, 0.78, 4.9, 6.2, 0.42)
    s.shapes.add_picture(str(QR), Inches(10.15), Inches(1.15), width=Inches(2.05), height=Inches(2.05))
    textbox(s, "Aponte a camera para abrir o GitHub", 9.65, 3.35, 3.0, 0.4, size=13, color="white", align=PP_ALIGN.CENTER)

    # Slide 2
    s = slides[1]
    textbox(s, "O que a demo mostra", 0.7, 0.55, 6.7, 0.5, size=30, bold=True)
    bullets(
        s,
        [
            "Criacao de um OCI Generative AI Agent em uma tenancy trial do zero.",
            "RAG com PDF do TDC Floripa 2026 para contexto geral do evento.",
            "Custom Tool HTTP para buscar agenda, speakers, trilhas e horarios.",
            "Agent Endpoint consumido por um canal externo: Telegram via backend Render.",
        ],
        0.85,
        1.45,
        6.1,
        4.3,
        size=19,
    )
    pill(s, "RAG = conhecimento nao estruturado", 7.45, 1.55, 4.6, 0.55, fill="neutral", line="ivy")
    pill(s, "Tool = dado estruturado e atualizavel", 7.45, 2.35, 4.6, 0.55, fill="neutral", line="sienna")
    pill(s, "Endpoint = integracao real", 7.45, 3.15, 4.6, 0.55, fill="neutral", line="ocean")
    pill(s, "Telegram = experiencia de canal", 7.45, 3.95, 4.6, 0.55, fill="neutral", line="rose")

    # Slide 3
    s = slides[2]
    textbox(s, "Arquitetura da solucao", 0.7, 0.45, 7.5, 0.5, size=30, bold=True)
    textbox(s, "Do canal externo ao Agent Endpoint, combinando RAG e Custom Tool.", 0.72, 0.95, 8.6, 0.35, size=15, color="muted")
    s.shapes.add_picture(str(ARCH), Inches(0.55), Inches(1.35), width=Inches(12.25), height=Inches(5.45))

    # Slide 4
    s = slides[3]
    textbox(s, "Roteiro da demonstracao", 0.7, 0.55, 7.3, 0.5, size=30, bold=True)
    bullets(
        s,
        [
            "1. Preparar OCI: compartment, policies, bucket, rede e subnet privada.",
            "2. Criar Knowledge Base e ingerir o PDF do evento.",
            "3. Criar Agent, adicionar RAG Tool e Custom Tool via OpenAPI.",
            "4. Testar perguntas: geral do evento, speaker, agenda e roteiro personalizado.",
            "5. Mostrar endpoint em canal externo com Telegram e Render.",
        ],
        0.85,
        1.35,
        6.4,
        4.8,
        size=18,
    )
    textbox(s, "Perguntas boas para testar", 7.75, 1.35, 4.3, 0.35, size=18, bold=True)
    bullets(
        s,
        [
            "O que sao as Jornadas TDC?",
            "Quais palestras a Livia Rodrigues vai fazer?",
            "Monte um roteiro para GenAI e arquitetura.",
        ],
        7.75,
        1.9,
        4.2,
        2.3,
        size=16,
        color="muted",
    )

    # Slide 5
    s = slides[4]
    textbox(s, "Material e proximos passos", 0.75, 0.8, 6.8, 0.6, size=32, bold=True, color="bark")
    textbox(
        s,
        "Use o repositorio para repetir o lab, atualizar os assets do TDC e adaptar a tool para outros eventos ou bases internas.",
        0.78,
        1.55,
        7.3,
        1.0,
        size=19,
        color="bark",
    )
    add_link(s, GITHUB_URL, 0.78, 3.1, 7.2, 0.45)
    s.shapes.add_picture(str(QR), Inches(9.75), Inches(1.1), width=Inches(2.15), height=Inches(2.15))
    textbox(s, "GitHub do lab", 9.6, 3.38, 2.5, 0.35, size=15, color="bark", align=PP_ALIGN.CENTER)

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    print(build())
